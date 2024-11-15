from flask import Flask, render_template
from flask import Flask, render_template, request, redirect, url_for, send_from_directory
import os
import re
from docx2txt import process as doc_process
from pptx import Presentation
from transformers import pipeline
import nltk

app = Flask(__name__)

app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['GENERATED_FOLDER'] = 'generated'

# Ensure necessary folders exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['GENERATED_FOLDER'], exist_ok=True)

# Route for index page
@app.route('/')
def index():
    return render_template('index.html')

# Route for topbar.html
@app.route('/topbar.html')
def topbar():
    return render_template('topbar.html')

# Route for sidebar.html
@app.route('/sidebar.html')
def sidebar():
    return render_template('sidebar.html')

# Route for maincontent.html
@app.route('/maincontent.html')
def maincontent():
    return render_template('maincontent.html')

# Route for aboutsystem.html
@app.route('/aboutsystem.html')
def aboutsystem():
    return render_template('aboutsystem.html')

# Route for aboutus.html
@app.route('/aboutus.html')
def aboutus():
    return render_template('aboutus.html')

# Define summarization pipeline here
def get_summarizer():
    return pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

summarizer = get_summarizer()


# New route for DOCX file upload and processing
@app.route('/upload', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        # Check if a file is provided
        if 'docx_file' not in request.files:
            return redirect(url_for('index'))
        
        file = request.files['docx_file']
        if file.filename == '':
            return redirect(url_for('index'))
        
        # Save uploaded file
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filepath)

        # Process and clean text from DOCX
        text = process_docx(filepath)
        cleaned_text = clean_unwanted_patterns(text)

        # Extract sections based on predefined markers
        section_markers = ["Course",
                           "Sem/AY",
                           "Module No.",
                           "Lesson Title",
                           "Description of the Lesson",
                           "Intended Learning Outcomes",
                           "Targets/ Objectives",
                           "Learning Guide Questions","Lecture Guide",
                           "Performance Task",
                           "Learning Resources"
                           ]
        sections = extract_sections_and_content(cleaned_text, section_markers)

        # Summarize the Lecture Guide section if available
        lecture_text = sections.get("Lecture Guide", "")
        if lecture_text:
            summary = summarizer(lecture_text, truncation=True, min_length=140, do_sample=False)
            summary_text = summary[0]['summary_text']
        else:
            summary_text = "No Lecture Guide content found for summarization."

        # Generate PPTX file from sections and summary
        ppt_file_path = os.path.join(app.config['GENERATED_FOLDER'], 'generated_presentation.pptx')
        create_ppt_with_word_limit(sections, summary_text, ppt_file_path, word_limit=50)

        # Redirect to download the generated PPTX
        return redirect(url_for('download_ppt', filename='generated_presentation.pptx'))

    return render_template('upload.html')

# Route to download the generated PPTX
@app.route('/download/<filename>')
def download_ppt(filename):
    return send_from_directory(app.config['GENERATED_FOLDER'], filename, as_attachment=True)

# DOCX processing functions
def process_docx(docx_path):
    return doc_process(docx_path)

def clean_unwanted_patterns(text, additional_patterns=None):
    unwanted_patterns = [
        r"Date\s*\n?.*\d{4}",
        r"Week\s*Duration\s*\n?.*\d+",
        r"Online\s*Activities\s*\(Synchronous/.*",
        r"Offline\s*Activities.*",
        r"Face to Face Activities \(Synchronous/.*",
        r"e-Learning/Self-Paced.*",
        r"Online\s*Discussion\s*via\s*Google\s*Meet.*",
        r"Student Learning Strategies.*",
        r"Asynchronous.*",
        r"LSPU SELF-PACED LEARNING MODULE: TECHNOLOGY FOR TEACHING AND LEARNING.*",
        r"The online discussion.*",
        r"Performance Task Date.*",
        r"You will be directed.*",
        r"The online discussion will happen.*"
    ]
    if additional_patterns:
        unwanted_patterns.extend(additional_patterns)
    
    for pattern in unwanted_patterns:
        text = re.sub(pattern, "", text, flags=re.MULTILINE | re.IGNORECASE)
    
    text = re.sub(r'\n\s*\n', '\n\n', text).strip()
    return text

def extract_sections_and_content(text, section_markers):
    sections = {}
    current_section = None
    buffer = []
    
    for line in text.splitlines():
        line = line.strip()
        if any(line.startswith(marker) for marker in section_markers):
            if current_section:
                sections[current_section] = "\n".join(buffer)
            current_section = line
            buffer = []
        elif current_section:
            buffer.append(line)
    
    if current_section:
        sections[current_section] = "\n".join(buffer)
    
    return sections

def create_ppt_with_word_limit(sections, summary_text, ppt_file_path, word_limit=50):
    presentation = Presentation()
    
    # Add title slide
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = sections.get("Lesson Title", "Untitled Document")

    # Add content slides
    for section, content in sections.items():
        if section != "Lesson Title":
            content_chunks = chunk_text(content, word_limit)
            for i, chunk in enumerate(content_chunks):
                slide_layout = presentation.slide_layouts[1]
                slide = presentation.slides.add_slide(slide_layout)
                slide.shapes.title.text = f"{section} (Part {i+1})" if i > 0 else section
                slide.placeholders[1].text = chunk

    # Add conclusion slide
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion"
    slide.placeholders[1].text = summary_text
    presentation.save(ppt_file_path)

def chunk_text(text, word_limit):
    words = text.split()
    chunks, chunk = [], []
    for word in words:
        chunk.append(word)
        if len(chunk) >= word_limit:
            chunks.append(" ".join(chunk))
            chunk = []
    if chunk:
        chunks.append(" ".join(chunk))
    return chunks


if __name__ == '__main__':
    app.run(debug=True)
